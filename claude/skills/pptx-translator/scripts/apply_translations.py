#!/usr/bin/env python3
"""抽出済みの all_runs.json と翻訳リスト(意味のあるテキストのみ)を突き合わせて
pptx にテキストだけを書き戻す。フォントサイズ・色・太字等の run.font 属性、
図形の位置・サイズ、罫線などの装飾には一切触れない。

translations.json は「空でない(strip()して空にならない)run」の数と同じ要素数の
JSON配列で、all_runs.json 内で空でない要素が出現する順序と1:1で対応させる。
空白/空文字のrunは自動的にスキップされ、そのまま変更されない。

使い方:
  python apply_translations.py <src.pptx> <all_runs.json> <translations.json> <dst.pptx>
"""
import json
import os
import sys
from pptx import Presentation


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:
            yield from iter_shapes(sh.shapes)


def main():
    if len(sys.argv) != 5:
        print("usage: apply_translations.py <src.pptx> <all_runs.json> <translations.json> <dst.pptx>",
              file=sys.stderr)
        sys.exit(1)
    src, all_runs_path, translations_path, dst = sys.argv[1:5]
    for path in (src, all_runs_path, translations_path):
        if not os.path.isfile(path):
            print(f"ERROR: ファイルが見つかりません: {path}", file=sys.stderr)
            sys.exit(1)

    original_all = json.load(open(all_runs_path, encoding="utf-8"))
    translations = json.load(open(translations_path, encoding="utf-8"))

    nonblank_count = sum(1 for t in original_all if t.strip() != "")
    if len(translations) != nonblank_count:
        print(f"ERROR: translations has {len(translations)} items, "
              f"but {nonblank_count} non-blank runs were found in {all_runs_path}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(src)

    idx = 0        # original_all 全体(空白含む)へのインデックス
    trans_pos = 0  # translations(意味のある行のみ)へのインデックス

    def apply_tf(tf):
        nonlocal idx, trans_pos
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text != original_all[idx]:
                    raise AssertionError(
                        f"drift at run index {idx}: actual={run.text!r} "
                        f"expected(from all_runs.json)={original_all[idx]!r}. "
                        "pptxの構造がall_runs.json生成時から変わっている可能性があります。"
                        "extract_texts.pyを再実行してください。")
                if original_all[idx].strip() != "":
                    run.text = translations[trans_pos]
                    trans_pos += 1
                idx += 1

    for slide in prs.slides:
        for sh in iter_shapes(slide.shapes):
            if sh.has_text_frame:
                apply_tf(sh.text_frame)
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        apply_tf(cell.text_frame)

    assert idx == len(original_all), f"applied {idx} of {len(original_all)} runs"
    assert trans_pos == len(translations), f"used {trans_pos} of {len(translations)} translations"

    prs.save(dst)
    print(f"saved {dst}: {idx} runs walked, {trans_pos} translations applied")


if __name__ == "__main__":
    main()
