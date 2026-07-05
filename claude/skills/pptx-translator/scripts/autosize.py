#!/usr/bin/env python3
"""テキストが図形(テキストフレーム/テーブルセル)の矩形からはみ出している場合に、
収まる最大フォントサイズを計算して縮小する。図形の位置・サイズ・色・罫線などの
装飾やレイアウトには一切触れない。既に矩形に収まっているテキストは変更しない。

文字幅は簡易的な近似(欧文フォント前提: 平均文字幅≒フォントサイズ×0.52、
行送り≒フォントサイズ×1.18)でワードラップ行数を見積もる。日本語などCJK
向けの文字幅とは異なるため、翻訳先が欧文言語の場合を主な対象とする。

使い方:
  python autosize.py <pptx> [--min-pt 9.0]

<pptx> を直接上書きする。実行前に必要であればコピーを取っておくこと。
"""
import argparse
import json
import os
import sys
from pptx import Presentation
from pptx.util import Pt

CHAR_WIDTH_RATIO = 0.52
LINE_HEIGHT_RATIO = 1.18
DEFAULT_L_IN = 91440
DEFAULT_R_IN = 91440
DEFAULT_T_IN = 45720
DEFAULT_B_IN = 45720


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:
            yield from iter_shapes(sh.shapes)


def wrapped_line_count(text, chars_per_line):
    if chars_per_line <= 0:
        return 999
    lines = 0
    for para_text in text.split("\n"):
        words = para_text.split(" ")
        if not words or para_text == "":
            lines += 1
            continue
        cur = 0
        line_count = 1
        for w in words:
            wl = len(w)
            add = wl if cur == 0 else wl + 1
            if cur + add > chars_per_line and cur > 0:
                line_count += 1
                cur = wl
            else:
                cur += add
        lines += line_count
    return lines


def required_lines(text, font_pt, avail_width_pt):
    chars_per_line = max(1, int(avail_width_pt / (font_pt * CHAR_WIDTH_RATIO)))
    return wrapped_line_count(text, chars_per_line)


def fits(text, font_pt, avail_width_pt, avail_height_pt):
    lines = required_lines(text, font_pt, avail_width_pt)
    needed_height = lines * font_pt * LINE_HEIGHT_RATIO
    return needed_height <= avail_height_pt


def max_fitting_size(text, current_pt, avail_width_pt, avail_height_pt, min_pt):
    if fits(text, current_pt, avail_width_pt, avail_height_pt):
        return current_pt
    lo, hi = min_pt, current_pt
    best = min_pt
    for _ in range(30):
        mid = (lo + hi) / 2
        if fits(text, mid, avail_width_pt, avail_height_pt):
            best = mid
            lo = mid
        else:
            hi = mid
        if hi - lo < 0.1:
            break
    return round(best * 2) / 2


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--min-pt", type=float, default=9.0)
    ap.add_argument("--report", default=None, help="変更点をJSONで書き出すパス(省略可)")
    args = ap.parse_args()
    if not os.path.isfile(args.pptx):
        print(f"ERROR: pptxファイルが見つかりません: {args.pptx}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(args.pptx)

    report = []
    changed = 0
    hit_min = []

    def process_text_frame(tf, width_emu, height_emu, si, shape_label):
        nonlocal changed
        text = tf.text
        if not text.strip():
            return

        l_in = tf.margin_left if tf.margin_left is not None else DEFAULT_L_IN
        r_in = tf.margin_right if tf.margin_right is not None else DEFAULT_R_IN
        t_in = tf.margin_top if tf.margin_top is not None else DEFAULT_T_IN
        b_in = tf.margin_bottom if tf.margin_bottom is not None else DEFAULT_B_IN

        avail_w_pt = (width_emu - l_in - r_in) / 12700
        avail_h_pt = (height_emu - t_in - b_in) / 12700
        if avail_w_pt <= 0 or avail_h_pt <= 0:
            return

        runs = [run for para in tf.paragraphs for run in para.runs if run.font.size is not None]
        if not runs:
            return

        current_max_pt = max(r.font.size.pt for r in runs)
        best = max_fitting_size(text, current_max_pt, avail_w_pt, avail_h_pt, args.min_pt)

        if best < current_max_pt - 0.1:
            scale = best / current_max_pt
            for run in runs:
                old = run.font.size.pt
                new = max(round(old * scale * 2) / 2, args.min_pt)
                run.font.size = Pt(new)
            changed += 1
            report.append({"slide": si, "shape": shape_label,
                            "before_max_pt": current_max_pt, "after_max_pt": best,
                            "text": text[:50]})
            if best <= args.min_pt + 0.01:
                hit_min.append((si, shape_label, text[:50]))

    for si, slide in enumerate(prs.slides):
        for sh in iter_shapes(slide.shapes):
            if sh.has_text_frame:
                try:
                    width_emu, height_emu = sh.width, sh.height
                except TypeError:
                    continue
                process_text_frame(sh.text_frame, width_emu, height_emu, si, f"shape_id={sh.shape_id}")
            if sh.has_table:
                tbl = sh.table
                for ri, row in enumerate(tbl.rows):
                    for ci, cell in enumerate(row.cells):
                        process_text_frame(cell.text_frame, tbl.columns[ci].width, row.height,
                                            si, f"table_id={sh.shape_id} cell[{ri}][{ci}]")

    prs.save(args.pptx)

    result = {"changed": changed, "hit_min_count": len(hit_min), "details": report, "hit_min": hit_min}
    if args.report:
        with open(args.report, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=1)
    print(f"changed={changed} hit_min={len(hit_min)}")
    if hit_min:
        print("以下は最小フォントサイズに達しても収まりきらない可能性があります(要目視確認):")
        for si, label, text in hit_min:
            print(f"  slide={si} {label} text={text!r}")


if __name__ == "__main__":
    main()
