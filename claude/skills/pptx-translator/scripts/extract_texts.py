#!/usr/bin/env python3
"""pptx内の全テキストランをdocument順に抽出する。

出力:
  <out_dir>/all_runs.json      - 全run(空白/空文字含む)をdocument順に並べたリスト
  <out_dir>/nonblank_pairs.txt - 翻訳対象(空でないテキストのみ)を連番付きで並べたレビュー用ファイル

使い方:
  python extract_texts.py <src.pptx> <out_dir>
"""
import json
import os
import sys
from pptx import Presentation


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # GROUP
            yield from iter_shapes(sh.shapes)


def collect_all_runs(prs):
    """text_frame と table cell の全runをdocument順(スライド→shape→段落→run)で集める。"""
    items = []
    for slide in prs.slides:
        for sh in iter_shapes(slide.shapes):
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        items.append(run.text)
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                items.append(run.text)
    return items


def main():
    if len(sys.argv) != 3:
        print("usage: extract_texts.py <src.pptx> <out_dir>", file=sys.stderr)
        sys.exit(1)
    src, out_dir = sys.argv[1], sys.argv[2]
    if not os.path.isfile(src):
        print(f"ERROR: pptxファイルが見つかりません: {src}", file=sys.stderr)
        sys.exit(1)
    if not os.path.isdir(out_dir):
        print(f"ERROR: 出力先ディレクトリが存在しません: {out_dir}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(src)
    items = collect_all_runs(prs)

    with open(f"{out_dir}/all_runs.json", "w", encoding="utf-8") as f:
        json.dump(items, f, ensure_ascii=False, indent=1)

    nonblank = [(i, t) for i, t in enumerate(items) if t.strip() != ""]
    lines = [f"{n}: orig_idx={i} {t!r}" for n, (i, t) in enumerate(nonblank)]
    with open(f"{out_dir}/nonblank_pairs.txt", "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"total_runs={len(items)} nonblank={len(nonblank)} blank={len(items) - len(nonblank)}")


if __name__ == "__main__":
    main()
